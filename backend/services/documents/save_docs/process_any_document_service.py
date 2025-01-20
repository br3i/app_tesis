import json
import re
import os
import uuid
import pptx
import pandas as pd
import pytesseract
import time
from docx import Document as DocxDocument
from PIL import Image
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PDFPlumberLoader
from langchain.schema import Document
from models.database import get_db
from services.embeddings.save_embedding_service import save_embeddings
from services.helpers.return_collection import return_collection
from services.documents.treat_docs.info_documents_service import get_info_document
from dotenv import load_dotenv

# Especifica la ruta al archivo .env
dotenv_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), '../../../.env')
load_dotenv(dotenv_path)

# Leer el contenido del archivo .env
CHUNK_SIZE = os.getenv('CHUNK_SIZE', '512')  # Establecer valor predeterminado
CHUNK_OVERLAP = os.getenv('CHUNK_OVERLAP', '40')  # Establecer valor predeterminado
LENGTH_FUNCTION = os.getenv('LENGTH_FUNCTION', 'len')  # Se asigna por defecto a 'len'
IS_SEPARATOR_REGEX = os.getenv('IS_SEPARATOR_REGEX', 'False')  # Predeterminado a 'False'

# Convertir las variables a los tipos correctos
try:
    # Convertir CHUNK_SIZE y CHUNK_OVERLAP a enteros
    CHUNK_SIZE = int(CHUNK_SIZE)
    CHUNK_OVERLAP = int(CHUNK_OVERLAP)
    
    # Convertir LENGTH_FUNCTION (especificamos la función que queremos usar)
    if LENGTH_FUNCTION == 'len':
        LENGTH_FUNCTION = len  # Asignamos la función len

    # Convertir IS_SEPARATOR_REGEX a booleano
    IS_SEPARATOR_REGEX = IS_SEPARATOR_REGEX.lower() in ['true', '1', 't', 'y', 'yes']
except ValueError as e:
    print(f"Error en la conversión de las variables del .env: {e}")
    # Asignar valores predeterminados si algo falla
    CHUNK_SIZE = 1024
    CHUNK_OVERLAP = 80
    LENGTH_FUNCTION = len
    IS_SEPARATOR_REGEX = False

# Ahora pasamos estas variables al text_splitter
print(f"[Valores del env] CHUNK_SIZE: {CHUNK_SIZE} ({type(CHUNK_SIZE)})")
print(f"[Valores del env] CHUNK_OVERLAP: {CHUNK_OVERLAP} ({type(CHUNK_OVERLAP)})")
print(f"[Valores del env] LENGTH_FUNCTION: {LENGTH_FUNCTION} ({type(LENGTH_FUNCTION)})")
print(f"[Valores del env] IS_SEPARATOR_REGEX: {IS_SEPARATOR_REGEX} ({type(IS_SEPARATOR_REGEX)})")

# Configuración del splitter de texto
# text_splitter = RecursiveCharacterTextSplitter(
#     chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP, length_function=LENGTH_FUNCTION, is_separator_regex=IS_SEPARATOR_REGEX
# )
text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP, length_function=LENGTH_FUNCTION, is_separator_regex=IS_SEPARATOR_REGEX
)


def ocr_from_pdf(file_path: str):
    """ Convertir PDF con imágenes a texto usando OCR """
    try:
        pages = convert_from_path(file_path)
        docs = []
        
        for page_number, page in enumerate(pages, start=1):
            text = pytesseract.image_to_string(page)
            if text.strip():  # Si se detectó texto en la imagen
                docs.append({
                    "page_content": text,
                    "metadata": {"page": page_number}
                })
        print(f"[ocr_from_pdf] Docs generados con OCR: {len(docs)}")
        return docs
    except Exception as e:
        print(f"Error al procesar OCR: {e}")
        return []

# Función para extraer el número de resolución
def extract_resolution_from_name(document_name):
    resolution_pattern = r'RESOLUCIÓN\s(\d+)\.CP\.(\d{4})'
    match = re.search(resolution_pattern, document_name)
    if match:
        return match.group(1)  # Devuelve el número de la resolución
    return None  # Si no se encuentra una resolución



def process_pdf(file_path: str, collection_name: str, id_document: int):
    #!!!!!!!!!!!!!!!!!!!!!REVISAR DOCUMENTO 27-2022 -> Error procesando el PDF: list index out of range
    print("\n\n--------------------------[PROCESS_PDF]--------------------------")
    print("[process_pdf] file_apath: ", file_path)
    print("[process_pdf] collection_name: ", collection_name)
    print("[process_pdf] id_document: ", id_document)
    """ Procesar documentos PDF """
    db = next(get_db())
    try:
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"El archivo {file_path} no existe.")

        resolution, articles_entities, copia, resolve, resolve_page = get_info_document(file_path)
        print(f"\n\n-resolution 1: \n{resolution}")
        if not resolution:
            resolution = [f'{os.path.splitext(os.path.basename(file_path))[0]}']
        # print("\n----------------------------------------------------------")
        # print(f"\n-resolution 2: \n{resolution}")
        # print(f"\n----------------------------------------------------------")
        # print(f"\n\n-articles_entities: {articles_entities} \n")
        # print(f"\n\n----------------------------------------------------------")
        # print(f"\n\n[proces_any_doc] resolve \n{resolve}")
        # print(f"\n\n----------------------------------------------------------")
        # print(f"\n\n----------------------------------------------------------")
        # print("\n\copia: ", copia)
        if resolve_page is None:
            resolve_page = 0
        # print("\n\nresolve_page: ", resolve_page)
        # print("\ntipo resolve_page: ", type(resolve_page))
        # print(f"\n\n----------------------------------------------------------")
        resolve_page = str(resolve_page)
        # print("\n\nresolve_page: ", resolve_page)
        # print("\ntipo resolve_page: ", type(resolve_page))
        # time.sleep(1000)

        documents = [Document(page_content=resolve)]
        # print(f"\n\n----------------------------------------------------------")
        # print(f"\n\nDocuments: {documents}")
        ##Generar los chunks a partir del contenido de 'resolve'
        text_chunks = text_splitter.split_documents(documents)
        # print(f"\n\n----------------------------------------------------------")
        # print(f"\n\nText chunks: {text_chunks}")
        
        ##Obtener el nombre de la resolución (de la variable 'resolution' que debes haber definido anteriormente)
        document_name = resolution  # Suponiendo que 'resolution' es la variable que contiene el nombre de la resolución
        # print(f"\n\n----------------------------------------------------------")
        # print(f"\n\nDocument name: {document_name}")

        # Procesar los artículos y asociarlos con su página
        considerations = []
        for article in articles_entities:
            consideration_metadata = {
                "consideration": article
            }
            considerations.append(consideration_metadata)
        # print(f"\n\n----------------------------------------------------------")
        # print(f"\n\nConsiderations: {considerations}")
        
        # Crear metadata base fuera del loop
        base_metadata = {
            "document_name": document_name,
            "file_path": file_path,
            "resolve_page": resolve_page,  # Asignamos la página de la resolución
            "collection_name": collection_name,
            "considerations": considerations,  # Asignamos las consideraciones a la metadata
            "copia": copia
        }
        # print(f"\n\n----------------------------------------------------------")
        # print(f"\n\nBase metadata: {base_metadata}")
        #time.sleep(10)

        # Crear la metadata para los fragmentos de 'resolve'
        for idx, chunk in enumerate(text_chunks):
            #time.sleep(5)
            try:
                print(f"[process_resolve_and_articles] Procesando fragmento {idx + 1}...")

                # Verificar tipo de chunk y sus atributos
                #print(f"[process_resolve_and_articles] Tipo de chunk: {type(chunk)}")
                #print(f"\n[process_resolve_and_articles] Contenido del chunk: {chunk}")  # Mostrar solo los primeros 100 caracteres

                fragment_id = str(uuid.uuid4())
                # print("[process_any_doc_service] fragment_id = ", fragment_id)
                
                # Crear metadata para el fragmento
                document_metadata = {
                    **base_metadata,
                    "uuid": fragment_id,
                    "chunk_index": str(idx),
                    "text": chunk.page_content, # Asignamos las consideraciones a la metadata
                }
                # print(f"\n\n----------------------------------------------------------")
                # print(f"\n\n[process_resolve_and_articles] Metadata para el fragmento: {document_metadata}")

                # Llamar a la función que guarda las embeddings
                collection = return_collection(collection_name)
                # print(f"\n\n----------------------------------------------------------")
                # print(f"\n\n[process_resolve_and_articles] Collection: {collection}")
                
                save_embeddings([chunk.page_content], collection, document_metadata, id_document, db)

            except Exception as e:
                print(f"Error procesando el fragmento {idx + 1}: {e}")
        return len(documents), len(text_chunks)
    except Exception as e:
        print(f"Error procesando el PDF: {e}")
        return 0, 0
    finally:
        db.close()


def process_word_document(file_path: str, collection_name: str):
    print("--------------------------[PROCESS_WORD_DOCUMENT]--------------------------")
    """ Procesar documentos Word (docx, doc) """
    try:
        doc = DocxDocument(file_path)
        print("[prc-wrd] valor de doc: ", doc)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        print("[prc-wrd] valor de full_text: ", full_text)
                
        # Dividir el texto en fragmentos
        text_chunks = text_splitter.split_documents(full_text)
        print(f"[prc-wrd] Fragmentos generados: {len(text_chunks)}")
        print(f"[prc-wrd] Valor de text_chunks {text_chunks}")
        
        # Guardar los fragmentos como embeddings
        for idx, chunk in enumerate(text_chunks):
            try:
                fragment_id = str(uuid.uuid4())
                document_metadata = {
                    "uuid": fragment_id,
                    "document_name": os.path.basename(file_path),
                    "chunk_index": str(idx),
                    "text": chunk,
                    "collection_name": collection_name
                }

                collection = return_collection(collection_name)
                
                save_embeddings([chunk], collection_name, document_metadata)
            except Exception as e:
                print(f"Error procesando el fragmento {idx + 1}: {e}")
        
        return len(full_text), len(text_chunks)
    except Exception as e:
        print(f"Error procesando el documento Word: {e}")
        return 0, 0


def process_text_document(file_path: str, collection_name: str):
    print("--------------------------[PROCESS_TEXT_DOCUMENT]--------------------------")
    """ Procesar documentos de texto (txt, rtf, odf) """
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        
        print("[prc-txt] Valor de content: ", content)
        
        # Dividir el texto en fragmentos
        text_chunks = text_splitter.split_documents([content])
        print(f"[prc-txt] Fragmentos generados: {len(text_chunks)}")
        print(f"[prc-txt] Valor de text_chunks {text_chunks}")
        
        # Guardar los fragmentos como embeddings
        for idx, chunk in enumerate(text_chunks):
            try:
                fragment_id = str(uuid.uuid4())
                document_metadata = {
                    "uuid": fragment_id,
                    "document_name": os.path.basename(file_path),
                    "chunk_index": str(idx),
                    "text": chunk,
                    "collection_name": collection_name
                }

                collection = return_collection(collection_name)
                
                save_embeddings([chunk], collection_name, document_metadata)
            except Exception as e:
                print(f"Error procesando el fragmento {idx + 1}: {e}")
        
        return 1, len(text_chunks)  # Devuelve 1 porque solo es un documento de texto
    except Exception as e:
        print(f"Error procesando el documento de texto: {e}")
        return 0, 0


def process_image(file_path: str, collection_name: str):
    print("--------------------------[PROCESS_IMAGE]--------------------------")
    """ Procesar imágenes con OCR utilizando pytesseract """
    try:
        image = Image.open(file_path)
        print(f"[prc-img] Valor de image {image}")
        text = pytesseract.image_to_string(image)
        print(f"[prc-img] Valor de text {text}")
        
        # Dividir el texto en fragmentos
        text_chunks = text_splitter.split_documents([text])
        print(f"[prc-img] Fragmentos generados: {len(text_chunks)}")
        print(f"[prc-img] Valor de text_chunks {text_chunks}")
        
        # Guardar los fragmentos como embeddings
        for idx, chunk in enumerate(text_chunks):
            try:
                fragment_id = str(uuid.uuid4())
                document_metadata = {
                    "uuid": fragment_id,
                    "document_name": os.path.basename(file_path),
                    "chunk_index": str(idx),
                    "text": chunk,
                    "collection_name": collection_name
                }

                collection = return_collection(collection_name)
                
                save_embeddings([chunk], collection_name, document_metadata)
            except Exception as e:
                print(f"Error procesando el fragmento {idx + 1}: {e}")
        
        return 1, len(text_chunks)  # Devuelve 1 ya que es una sola imagen
    except Exception as e:
        print(f"Error procesando la imagen: {e}")
        return 0, 0


def process_ppt(file_path: str, collection_name: str):
    print("--------------------------[PROCESS_PPT]--------------------------")
    """ Procesar presentaciones de PowerPoint (ppt, pptx) """
    try:
        presentation = pptx.Presentation(file_path)
        print(f"[prc-ppt] Valor de presentation {presentation}")
        full_text = []
        
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        
        print(f"[prc-ppt] Valor de full_text {full_text}")
        
        # Dividir el texto en fragmentos
        text_chunks = text_splitter.split_documents(full_text)
        print(f"[prc-ppt] Fragmentos generados: {len(text_chunks)}")
        print(f"[prc-ppt] Valor de text_chunks {text_chunks}")
        
        # Guardar los fragmentos como embeddings
        for idx, chunk in enumerate(text_chunks):
            try:
                fragment_id = str(uuid.uuid4())
                document_metadata = {
                    "uuid": fragment_id,
                    "document_name": os.path.basename(file_path),
                    "chunk_index": str(idx),
                    "text": chunk,
                    "collection_name": collection_name
                }

                collection = return_collection(collection_name)
                
                save_embeddings([chunk], collection_name, document_metadata)
            except Exception as e:
                print(f"Error procesando el fragmento {idx + 1}: {e}")
        
        return len(full_text), len(text_chunks)
    except Exception as e:
        print(f"Error procesando la presentación PPT: {e}")
        return 0, 0


def process_excel(file_path: str, collection_name: str):
    print("--------------------------[PROCESS_EXCEL]--------------------------")
    """ Procesar documentos Excel (xls, xlsx) """
    try:
        df = pd.read_excel(file_path)
        print(f"[prc-excel] Valor de df {df}")
        full_text = df.to_string(index=False, header=False)  # Convertir todo a texto
        print(f"[prc-excel] Valor de full_text {full_text}")
        
        # Dividir el texto en fragmentos
        text_chunks = text_splitter.split_documents([full_text])
        print(f"[prc-excel] Fragmentos generados: {len(text_chunks)}")
        print(f"[prc-excel] Valor de text_chunks {text_chunks}")
        
        # Guardar los fragmentos como embeddings
        for idx, chunk in enumerate(text_chunks):
            try:
                fragment_id = str(uuid.uuid4())
                document_metadata = {
                    "uuid": fragment_id,
                    "document_name": os.path.basename(file_path),
                    "chunk_index": str(idx),
                    "text": chunk,
                    "collection_name": collection_name
                }

                collection = return_collection(collection_name)
                
                save_embeddings([chunk], collection_name, document_metadata)
            except Exception as e:
                print(f"Error procesando el fragmento {idx + 1}: {e}")
        
        return 1, len(text_chunks)  # Devuelve 1 ya que es un solo archivo Excel
    except Exception as e:
        print(f"Error procesando el documento Excel: {e}")
        return 0, 0
